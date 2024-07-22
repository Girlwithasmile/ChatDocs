from flask import Flask, request, render_template, jsonify
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain.schema import Document
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

app = Flask(__name__)

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
    return text

def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    return text

def get_docx_text(docx_docs):
    text = ""
    for docx_file in docx_docs:
        doc = DocxDocument(docx_file)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
    return text

def get_text_from_files(files):
    text = ""
    pdf_files = [f for f in files if f.filename.endswith('.pdf')]
    ppt_files = [f for f in files if f.filename.endswith('.pptx')]
    docx_files = [f for f in files if f.filename.endswith('.docx')]

    if pdf_files:
        text += get_pdf_text(pdf_files)
    if ppt_files:
        text += get_ppt_text(ppt_files)
    if docx_files:
        text += get_docx_text(docx_files)
    
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    if not text_chunks:
        raise ValueError("Text chunks are empty")
    
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    try:
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("faiss_index")
    except Exception as e:
        raise Exception(f"Error creating FAISS index: {e}")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    try:
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    except Exception as e:
        return str(e)
    
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def summarize_text(text):
    prompt_template = """
    Summarize the following text in detail:\n\n
    Text:\n {context}\n

    Summary:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    document = Document(page_content=text)
    response = chain({"input_documents": [document]}, return_only_outputs=True)
    return response["output_text"]

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_files():
    if 'files' not in request.files:
        return jsonify({'error': 'No files part'})
    files = request.files.getlist('files')
    raw_text = get_text_from_files(files)
    if raw_text:
        text_chunks = get_text_chunks(raw_text)
        if text_chunks:
            try:
                get_vector_store(text_chunks)
                # Summarize the text after processing files
                summary = summarize_text(raw_text)
                return jsonify({'message': 'Files processed and index created successfully', 'summary': summary})
            except Exception as e:
                return jsonify({'error': str(e)})
        else:
            return jsonify({'error': 'No text chunks were generated.'})
    else:
        return jsonify({'error': 'No text was extracted from the files.'})

@app.route('/summarize', methods=['POST'])
def summarize():
    if 'files' not in request.files:
        return jsonify({'error': 'No files part'})
    files = request.files.getlist('files')
    raw_text = get_text_from_files(files)
    if raw_text:
        try:
            summary = summarize_text(raw_text)
            return jsonify({'summary': summary})
        except Exception as e:
            return jsonify({'error': str(e)})
    else:
        return jsonify({'error': 'No text was extracted from the files.'})

@app.route('/ask', methods=['POST'])
def ask():
    question = request.form.get('question')
    if not question:
        return jsonify({'error': 'No question provided'})
    
    try:
        answer = user_input(question)
        return jsonify({'answer': answer})
    except Exception as e:
        return jsonify({'error': str(e)})

if __name__ == '__main__':
    app.run(debug=True)
